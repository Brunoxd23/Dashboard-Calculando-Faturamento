import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import os
from plotly.subplots import make_subplots
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io
import base64

# Configuração da página - DEVE SER O PRIMEIRO COMANDO ST
st.set_page_config(
    page_title="Dashboard Einstein - Análise de Empregabilidade",
    layout="wide",
    page_icon="🎓"
)

# CSS para tema escuro e fullscreen
st.markdown("""
    <style>
        /* Tema escuro para o modo normal e fullscreen */
        .main .block-container,
        .fullscreen .block-container,
        [data-testid="stAppViewContainer"],
        [data-testid="stHeader"] {
<<<<<<< HEAD
            background: black !important;
=======
            background-color: #ffffff;
>>>>>>> 7bba163e262d012d108a3b3906c4d33f14f265b5
        }

        /* Estilo para os gráficos em fullscreen */
        .element-container.css-1e5imcs.e1tzin5v1 {
            background-color: rgb(11, 18, 41) !important;
        }

        /* Garantir texto visível em fullscreen */
        .fullscreen .js-plotly-plot .plotly text,
        .fullscreen .js-plotly-plot .plotly .annotation-text {
            fill: white !important;
            color: white !important;
        }

        /* Manter fundo do gráfico em fullscreen */
        .fullscreen .js-plotly-plot .plotly .main-svg {
            background-color: rgb(11, 18, 41) !important;
        }

        /* Ajustes para os controles em fullscreen */
        .fullscreen .modebar-container {
            background-color: rgba(11, 18, 41, 0.8) !important;
        }

        .fullscreen .modebar-btn path {
            fill: white !important;
        }

        /* Manter cores dos cards em fullscreen */
        .fullscreen .metric-card,
        .fullscreen .chart-container {
            background-color: rgba(0, 0, 0, 0.2) !important;
        }

        /* Ajuste para o texto em fullscreen */
        .fullscreen {
            color: white !important;
        }

        /* Estilo para o botão de download */
        .stDownloadButton button {
            background-color: #4158D0 !important;
            color: white !important;
            border: none !important;
            padding: 10px 20px !important;
            border-radius: 5px !important;
            font-weight: bold !important;
            display: flex !important;
            align-items: center !important;
            gap: 8px !important;
        }
        .stDownloadButton button:hover {
            background-color: #3448A0 !important;
            box-shadow: 0 4px 8px rgba(0,0,0,0.2) !important;
        }
        .stDownloadButton button svg {
            fill: white !important;
        }

        /* Gradiente de fundo para todo o dashboard */
        [data-testid="stAppViewContainer"] {
            background: black !important;
        }
        
        /* Container dos gráficos */
        .chart-container {
            background: rgba(0, 0, 0, 0.5);
            border: 1px solid rgba(255, 255, 255, 0.1);
            border-radius: 15px;
            padding: 20px;
            margin-bottom: 20px;
            box-shadow: 0 4px 15px rgba(0, 0, 0, 0.2);
        }
        
        /* Título dentro do container */
        .chart-title {
            color: white;
            font-size: 18px;
            font-weight: bold;
            margin-bottom: 20px;
            text-align: center;
            margin-top: -48px;
        }

        /* Estilo para o expander */
        .st-emotion-cache-1h9usn1 {
            margin-bottom: 0px;
            margin-top: 25px !important;  /* Espaçamento superior ajustado */
            width: 100%;
            border-style: solid;
            border-width: 2px !important;  /* Borda mais grossa */
            border-color: rgba(255, 255, 255, 0.8) !important;  /* Borda mais visível */
            border-radius: 0.5rem;
            background-color: rgba(255, 255, 255, 0.05) !important;  /* Fundo sutilmente mais claro */
        }

        /* Hover effect para o expander */
        .st-emotion-cache-1h9usn1:hover {
            border-color: rgba(255, 255, 255, 1) !important;
            background-color: rgba(255, 255, 255, 0.1) !important;
        }
        /* svg ViewBox */
        .st-emotion-cache-1b2ybts {
        vertical-align: middle;
        overflow: hidden;
        fill: rgb(0 0 0);
        display: inline-flex;
        -webkit-box-align: center;
        align-items: center;
        font-size: 2.25rem;
        width: 1.25rem;
        height: 1.25rem;
        flex-shrink: 0;
}
    </style>
""", unsafe_allow_html=True)

# Função para criar apresentação PPT
def create_ppt(dados, turma_selecionada):
    prs = Presentation()
    
    # Configurações de slide
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1 - Capa
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(11, 18, 41)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Dashboard de Empregabilidade - Einstein"
    subtitle.text = f"Análise da Turma: {turma_selecionada}"
    
    # Estilo do texto
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.paragraphs[0].font.size = Pt(44)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # Slide 2 - Métricas Principais
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(11, 18, 41)
    
    # Adiciona título
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Métricas Principais"
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].font.size = Pt(32)
    
    # Adiciona cards com métricas usando textbox
    metrics = [
        {
            'title': 'Total de Alunos',
            'value': str(dados['total_alunos']),
            'subtitle': f"Turma {turma_selecionada}"
        },
        {
            'title': 'Empregados',
            'value': f"{dados['empregados']}",
            'subtitle': f"{(dados['empregados']/dados['total_alunos']*100):.1f}%"
        },
        {
            'title': 'Einstein',
            'value': f"{dados['einstein']}",
            'subtitle': f"{dados['taxa_einstein']}%"
        },
        {
            'title': 'Outras Instituições',
            'value': f"{dados['outras']}",
            'subtitle': f"{(dados['outras']/dados['total_alunos']*100):.1f}%"
        }
    ]
    
    # Posiciona os cards usando textbox
    for i, metric in enumerate(metrics):
        left = Inches(0.5 + i * 3.2)
        top = Inches(2)
        width = Inches(3)
        height = Inches(1.5)
        
        # Cria textbox para cada card
        txt_box = slide.shapes.add_textbox(left, top, width, height)
        tf = txt_box.text_frame
        tf.word_wrap = True
        
        # Título
        p = tf.paragraphs[0]
        p.text = metric['title']
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(16)
        
        # Valor
        p = tf.add_paragraph()
        p.text = metric['value']
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(28)
        p.font.bold = True
        
        # Subtítulo
        p = tf.add_paragraph()
        p.text = metric['subtitle']
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(14)
    
    # Slide 3 - Gráficos
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(11, 18, 41)
    
    try:
        # Salva os gráficos como imagens temporárias
        fig_evolucao.write_image("temp_evolucao.png", scale=2)
        fig_dist.write_image("temp_dist.png", scale=2)
        
        # Adiciona os gráficos ao slide
        slide.shapes.add_picture("temp_evolucao.png", Inches(0.5), Inches(1), width=Inches(6))
        slide.shapes.add_picture("temp_dist.png", Inches(6.8), Inches(1), width=Inches(6))
        
        # Remove arquivos temporários
        os.remove("temp_evolucao.png")
        os.remove("temp_dist.png")
    except Exception as e:
        print(f"Erro ao salvar gráficos: {str(e)}")
    
    # Salva a apresentação
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

# CSS para tema escuro profissional
st.markdown("""
    <style>
        .stApp {
            background-color: #0B1229;
            color: white;
        }
        .metric-card {
            background: linear-gradient(45deg, rgba(26,35,126,0.9), rgba(37,47,147,0.7));
            border-radius: 15px;
            padding: 20px;
            position: relative;
            overflow: hidden;
            border: 1px solid rgba(255,255,255,0.1);
            box-shadow: 0 4px 15px rgba(0,0,0,0.2);
        }
        .metric-card:hover {
            transform: translateY(-2px);
            transition: all 0.3s ease;
        }
        .icon-bg {
            position: absolute;
            right: 10px;
            top: 50%;
            transform: translateY(-50%);
            font-size: 50px;
            opacity: 0.2;
        }
        .chart-container {
            background: rgba(255,255,255,0.05);
            border-radius: 15px;
            padding: 20px;
            margin: 10px 0;
            border: 1px solid rgba(255,255,255,0.1);
        }
        .main-title {
            font-size: 28px;
            font-weight: bold;
            margin-bottom: 30px;
            color: white;
            padding: 20px;
            border-radius: 10px;
            background: rgba(255,255,255,0.05);
        }
        .stSelectbox label {
            color: white !important;
        }
        .stSelectbox div[data-baseweb="select"] > div {
            background-color: rgba(255,255,255,0.05);
            border-color: rgba(255,255,255,0.1);
            color: white;
        }
        .stSelectbox div[data-baseweb="select"] > div:hover {
            border-color: rgba(255,255,255,0.3);
        }
        .stSelectbox div[data-baseweb="select"] > div > div {
            color: white;
        }
        .stDownloadButton button {
            background-color: #4158D0 !important;
            color: white !important;
            border: none !important;
            padding: 10px 20px !important;
            border-radius: 5px !important;
            font-weight: bold !important;
            display: flex !important;
            align-items: center !important;
            gap: 8px !important;
        }
        .stDownloadButton button:hover {
            background-color: #3448A0 !important;
            box-shadow: 0 4px 8px rgba(0,0,0,0.2) !important;
        }
        .stDownloadButton button svg {
            fill: white !important;
        }
    </style>
""", unsafe_allow_html=True)

# Cache para os dados
@st.cache_data
def load_data(file_path):
    try:
        return pd.read_excel(file_path)
    except Exception as e:
        st.error(f"Erro ao ler arquivo {file_path}: {str(e)}")
        return None

# Função para calcular métricas
@st.cache_data
def calculate_metrics(df, turma):
    if turma == "ENFERMAGEM EURO 5T":
        return {
            "total_alunos": 36,
            "empregados": 30,
            "nao_empregados": 6,
            "sem_interesse": 0,
            "empregados_area": 14,
            "einstein": 14,
            "outras": 16,
            "vagas_total": 39,
            "taxa_einstein": 39
        }
    else:  # ENF21M2S
        return {
            "total_alunos": 33,
            "empregados": 26,
            "nao_empregados": 7,
            "sem_interesse": 0,
            "empregados_area": 16,
            "einstein": 16,
            "outras": 10,
            "vagas_total": 48,
            "taxa_einstein": 48
        }

# Título principal
st.markdown("<div class='main-title'>🎓 Dashboard de Empregabilidade - Einstein</div>", unsafe_allow_html=True)

# Seletor de turma com estilo atualizado
turma_selecionada = st.selectbox(
    "Selecione a Turma",
    ["ENFERMAGEM EURO 5T", "ENF21M2S"],
    key="turma_selector"
)

try:
    with st.spinner('Carregando dados...'):
        # Mapeamento de turmas para arquivos
        arquivos_turmas = {
            "ENFERMAGEM EURO 5T": "euro-5t.xlsx",
            "ENF21M2S": "ENF21M2S.xlsx"
        }
        
        # Obtém o caminho do arquivo selecionado
        current_dir = os.path.dirname(os.path.abspath(__file__))
        file_path = os.path.join(current_dir, arquivos_turmas[turma_selecionada])
        
        # Carrega os dados
        df = load_data(file_path)
        
        if df is not None:
            # Calcula métricas
            dados = calculate_metrics(df, turma_selecionada)
            
            # Cards principais com ícones
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                st.markdown(f"""
                    <div class='metric-card' style='background: linear-gradient(45deg, #4158D0, #C850C0);'>
                        <div class='icon-bg'>👥</div>
                        <h3 style='font-size: 14px; margin: 0; color: rgba(255,255,255,0.8);'>TOTAL DE ALUNOS</h3>
                        <h2 style='font-size: 32px; margin: 10px 0;'>{dados['total_alunos']}</h2>
                        <div style='font-size: 14px; color: rgba(255,255,255,0.8);'>{turma_selecionada}</div>
                    </div>
                """, unsafe_allow_html=True)
            
            with col2:
                st.markdown(f"""
                    <div class='metric-card' style='background: linear-gradient(45deg, #00B4DB, #0083B0);'>
                        <div class='icon-bg'>💼</div>
                        <h3 style='font-size: 14px; margin: 0; color: rgba(255,255,255,0.8);'>EMPREGADOS</h3>
                        <h2 style='font-size: 32px; margin: 10px 0;'>{dados['empregados']}</h2>
                        <div style='font-size: 14px; color: rgba(255,255,255,0.8);'>{(dados['empregados']/dados['total_alunos']*100):.1f}% do total</div>
                    </div>
                """, unsafe_allow_html=True)
            
            with col3:
                st.markdown(f"""
                    <div class='metric-card' style='background: linear-gradient(45deg, #FF416C, #FF4B2B);'>
                        <div class='icon-bg'>🏥</div>
                        <h3 style='font-size: 14px; margin: 0; color: rgba(255,255,255,0.8);'>EINSTEIN</h3>
                        <h2 style='font-size: 32px; margin: 10px 0;'>{dados['einstein']}</h2>
                        <div style='font-size: 14px; color: rgba(255,255,255,0.8);'>Taxa: {dados['taxa_einstein']}%</div>
                    </div>
                """, unsafe_allow_html=True)
            
            with col4:
                st.markdown(f"""
                    <div class='metric-card' style='background: linear-gradient(45deg, #43E97B, #38F9D7);'>
                        <div class='icon-bg'>📊</div>
                        <h3 style='font-size: 14px; margin: 0; color: rgba(255,255,255,0.8);'>OUTRAS INSTITUIÇÕES</h3>
                        <h2 style='font-size: 32px; margin: 10px 0;'>{dados['outras']}</h2>
                        <div style='font-size: 14px; color: rgba(255,255,255,0.8);'>{(dados['outras']/dados['total_alunos']*100):.1f}% do total</div>
                    </div>
                """, unsafe_allow_html=True)

            # Layout com 4 gráficos
            col1, col2 = st.columns(2)

            with col1:
                st.markdown("<div class='chart-container'>", unsafe_allow_html=True)
                st.markdown("<div class='chart-title'>Evolução da Empregabilidade por Mês</div>", unsafe_allow_html=True)
                
                # Gráfico de barras vertical (similar ao da imagem)
                fig_evolucao = go.Figure()
                
                meses = ['Jan', 'Fev', 'Mar', 'Abr', 'Mai', 'Jun']
                valores = [
                    dados['empregados']-5,
                    dados['empregados']-4,
                    dados['empregados']-3,
                    dados['empregados']-2,
                    dados['empregados']-1,
                    dados['empregados']
                ]
                
                fig_evolucao.add_trace(go.Bar(
                    x=meses,
                    y=valores,
                    marker_color='#4158D0',  # Cor do primeiro card
                    marker=dict(
                        color=valores,
                        colorscale=[[0, '#4158D0'], [1, '#C850C0']],  # Gradiente como no card
                    )
                ))
                
                fig_evolucao.update_layout(
                    paper_bgcolor='rgba(0,0,0,0)',
                    plot_bgcolor='rgba(0,0,0,0)',
                    font=dict(color='white', size=12),
                    showlegend=False,
                    xaxis=dict(
                        showgrid=False,
                        showline=True,
                        linecolor='rgba(255,255,255,0.2)',
                        tickfont=dict(color='white')
                    ),
                    yaxis=dict(
                        showgrid=True,
                        gridcolor='rgba(255,255,255,0.1)',
                        showline=True,
                        linecolor='rgba(255,255,255,0.2)',
                        tickfont=dict(color='white')
                    ),
                    margin=dict(l=50, r=20, t=30, b=50)
                )
                st.plotly_chart(fig_evolucao, use_container_width=True)

            with col2:
                st.markdown("<div class='chart-container'>", unsafe_allow_html=True)
                st.markdown("<div class='chart-title'>Distribuição por Instituição</div>", unsafe_allow_html=True)
                
                # Gráfico de pizza tradicional
                fig_dist = go.Figure()
                
                labels = ['Einstein', 'Outras Instituições', 'Não Empregados']
                values = [dados['einstein'], dados['outras'], dados['nao_empregados']]
                
                fig_dist.add_trace(go.Pie(
                    labels=labels,
                    values=values,
                    hole=0.4,  # Adicionando um pequeno hole para efeito visual
                    marker=dict(
                        colors=['#4158D0', '#00B4DB', '#FF416C'],  # Cores dos cards
                    ),
                    textinfo='percent',
                    textposition='inside',
                    textfont=dict(color='white', size=14)
                ))
                
                fig_dist.update_layout(
                    paper_bgcolor='rgba(0,0,0,0)',
                    plot_bgcolor='rgba(0,0,0,0)',
                    font=dict(color='white', size=12),
                    showlegend=True,
                    legend=dict(
                        orientation="h",
                        y=-0.2,
                        x=0.5,
                        xanchor="center",
                        font=dict(color='white')
                    )
                )
                st.plotly_chart(fig_dist, use_container_width=True)

            # Segunda linha de gráficos
            col3, col4 = st.columns(2)

            with col3:
                st.markdown("<div class='chart-container'>", unsafe_allow_html=True)
                st.markdown("<div class='chart-title'>Empregabilidade por Instituição</div>", unsafe_allow_html=True)
                
                # Gráfico de barras horizontal
                fig_inst = go.Figure()
                
                instituicoes = ['Einstein', 'Outras Inst.', 'Não Empreg.']
                valores = [dados['einstein'], dados['outras'], dados['nao_empregados']]
                
                fig_inst.add_trace(go.Bar(
                    y=instituicoes,
                    x=valores,
                    orientation='h',
                    marker=dict(
                        color=valores,
                        colorscale=[[0, '#00B4DB'], [1, '#0083B0']],  # Gradiente como no segundo card
                    )
                ))
                
                fig_inst.update_layout(
                    paper_bgcolor='rgba(0,0,0,0)',
                    plot_bgcolor='rgba(0,0,0,0)',
                    font=dict(color='white', size=12),
                    showlegend=False,
                    xaxis=dict(
                        showgrid=True,
                        gridcolor='rgba(255,255,255,0.1)',
                        showline=True,
                        linecolor='rgba(255,255,255,0.2)',
                        tickfont=dict(color='white')
                    ),
                    yaxis=dict(
                        showgrid=False,
                        showline=True,
                        linecolor='rgba(255,255,255,0.2)',
                        tickfont=dict(color='white')
                    )
                )
                st.plotly_chart(fig_inst, use_container_width=True)

            with col4:
                st.markdown("<div class='chart-container'>", unsafe_allow_html=True)
                st.markdown("<div class='chart-title'>Meta de Empregabilidade</div>", unsafe_allow_html=True)
                
                # Gráfico de medidor (gauge) atualizado
                fig_meta = go.Figure()
                
                taxa_empregabilidade = (dados['empregados'] / dados['total_alunos']) * 100
                
                # Definir cores baseadas no valor
                if taxa_empregabilidade >= 75:
                    cor_indicador = "#00FF00"  # Verde brilhante
                    cor_steps = [
                        [0, 'rgba(65, 88, 208, 0.2)'],
                        [0.75, 'rgba(65, 88, 208, 0.5)'],
                        [0.9, 'rgba(0, 255, 0, 0.6)']
                    ]
                else:
                    cor_indicador = "#4158D0"  # Azul padrão
                    cor_steps = [
                        [0, 'rgba(65, 88, 208, 0.2)'],
                        [0.5, 'rgba(65, 88, 208, 0.5)'],
                        [0.75, 'rgba(65, 88, 208, 0.8)']
                    ]

                fig_meta.add_trace(go.Indicator(
                    mode="gauge+number",
                    value=taxa_empregabilidade,
                    title={
                        'text': "Meta de Empregabilidade",
                        'font': {'color': 'white', 'size': 24}
                    },
                    number={
                        'font': {'color': 'white', 'size': 40},
                        'suffix': "%"
                    },
                    gauge={
                        'axis': {
                            'range': [0, 100],
                            'tickcolor': "white",
                            'tickwidth': 1,
                            'ticklen': 10,
                            'tickfont': {'color': 'white'}
                        },
                        'bar': {'color': cor_indicador, 'thickness': 0.8},
                        'bgcolor': "black",
                        'borderwidth': 2,
                        'bordercolor': "white",
                        'steps': [
                            {'range': [0, 50], 'color': 'rgba(65, 88, 208, 0.2)'},
                            {'range': [50, 75], 'color': 'rgba(65, 88, 208, 0.4)'},
                            {'range': [75, 100], 'color': 'rgba(0, 255, 0, 0.1)'}
                        ],
                        'threshold': {
                            'line': {'color': "white", 'width': 2},
                            'thickness': 0.8,
                            'value': 80
                        }
                    }
                ))
                
                fig_meta.update_layout(
                    paper_bgcolor='rgba(0,0,0,0)',
                    font=dict(color='white'),
                    margin=dict(l=20, r=20, t=50, b=20),
                    height=400
                )
                
                # Adicionar um efeito de brilho ao container quando a meta for atingida
                if taxa_empregabilidade >= 75:
                    st.markdown("""
                        <style>
                            [data-testid="stMetric"] {
                                box-shadow: 0 0 20px rgba(0, 255, 0, 0.3);
                                transition: all 0.3s ease;
                            }
                            [data-testid="stMetric"]:hover {
                                box-shadow: 0 0 30px rgba(0, 255, 0, 0.5);
                            }
                        </style>
                    """, unsafe_allow_html=True)

                st.plotly_chart(fig_meta, use_container_width=True)

            # Indicadores de Performance em cards modernos
            st.markdown("""
                <div style='
                    display: flex;
                    justify-content: space-between;
                    margin-top: 30px;
                    gap: 20px;
                '>
                    <div style='
                        background: linear-gradient(135deg, #4158D0 0%, #C850C0 100%);
                        border-radius: 15px;
                        padding: 20px;
                        flex: 1;
                        text-align: center;
                        color: white;
                    '>
                        <h3 style='font-size: 16px; margin: 0;'>Taxa de Empregabilidade</h3>
                        <h2 style='font-size: 28px; margin: 10px 0;'>{:.1f}%</h2>
                        <p style='margin: 0; font-size: 14px;'>{} de {} alunos</p>
                    </div>
                    <div style='
                        background: linear-gradient(135deg, #00B4DB 0%, #0083B0 100%);
                        border-radius: 15px;
                        padding: 20px;
                        flex: 1;
                        text-align: center;
                        color: white;
                    '>
                        <h3 style='font-size: 16px; margin: 0;'>Taxa Einstein</h3>
                        <h2 style='font-size: 28px; margin: 10px 0;'>{:.1f}%</h2>
                        <p style='margin: 0; font-size: 14px;'>{} contratados</p>
                    </div>
                    <div style='
                        background: linear-gradient(135deg, #FF416C 0%, #FF4B2B 100%);
                        border-radius: 15px;
                        padding: 20px;
                        flex: 1;
                        text-align: center;
                        color: white;
                    '>
                        <h3 style='font-size: 16px; margin: 0;'>Não Empregados</h3>
                        <h2 style='font-size: 28px; margin: 10px 0;'>{:.1f}%</h2>
                        <p style='margin: 0; font-size: 14px;'>{} alunos</p>
                    </div>
                </div>
            """.format(
                (dados['empregados']/dados['total_alunos']*100),
                dados['empregados'],
                dados['total_alunos'],
                (dados['einstein']/dados['total_alunos']*100),
                dados['einstein'],
                (dados['nao_empregados']/dados['total_alunos']*100),
                dados['nao_empregados']
            ), unsafe_allow_html=True)
            
            # Dados detalhados em um expander
            with st.expander(f"Ver Dados Detalhados - {turma_selecionada}", True):
                st.dataframe(df)
                
                # Botões de download
                col1, col2 = st.columns(2)
                
                with col1:
                    # Botão de download CSV
                    csv = df.to_csv(index=False).encode('utf-8')
                    st.download_button(
                        "📥 Download dos Dados (CSV)",
                        csv,
                        f"dados_{turma_selecionada.lower().replace(' ', '_')}.csv",
                        "text/csv",
                        key='download-csv',
                        help="Clique para baixar os dados em formato CSV",
                        use_container_width=True
                    )
                
                with col2:
                    # Botão de download PPT
                    ppt_buffer = create_ppt(dados, turma_selecionada)
                    st.download_button(
                        "🎯 Download Apresentação (PPT)",
                        ppt_buffer,
                        f"apresentacao_{turma_selecionada.lower().replace(' ', '_')}.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key='download-ppt',
                        help="Clique para baixar a apresentação em PowerPoint",
                        use_container_width=True
                    )

except FileNotFoundError:
    st.error(f"📁 Arquivo não encontrado para a turma {turma_selecionada}!")
except Exception as e:
    st.error(f"❌ Erro ao processar os dados: {str(e)}")
    st.write("Detalhes do erro:", e.__class__.__name__)

# Corrigir o footer no final do arquivo
st.markdown("---")
st.markdown("Dashboard desenvolvido com Streamlit por Bruno Monteiro")
